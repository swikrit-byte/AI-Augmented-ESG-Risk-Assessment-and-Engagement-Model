"""
ESGIntel — PowerPoint Engagement Deck Generator
================================================
Generates McKinsey-style .pptx engagement decks.

Public contract:
    build_deck(deck_key: str, out_dir: str) -> str
        deck_key in {'board', 'investor'}
        Returns absolute path to the generated .pptx file.
"""

import os
import sys
import tempfile
import datetime
import textwrap

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Import company data (sibling module)
# ---------------------------------------------------------------------------
import company_data as D

# ---------------------------------------------------------------------------
# Slide dimensions (16:9)
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Colour helpers
# ---------------------------------------------------------------------------

def rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


INK     = rgb(D.BRAND["ink"])
ACCENT  = rgb(D.BRAND["accent"])
ACCENT_LT = rgb(D.BRAND["accent_lt"])
CORAL   = rgb(D.BRAND["coral"])
CORAL_LT = rgb(D.BRAND["coral_lt"])
SLATE   = rgb(D.BRAND["slate"])
MUTED   = rgb(D.BRAND["muted"])
RULE    = rgb(D.BRAND["rule"])
RED     = rgb(D.BRAND["red"])
AMBER   = rgb(D.BRAND["amber"])
GREEN   = rgb(D.BRAND["green"])
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
PAGE_BG = rgb(D.BRAND["page_bg"])
BAND_BG = rgb(D.BRAND["band_bg"])
CREAM   = rgb(D.BRAND["cream"])
ROW_ALT = rgb(D.BRAND["row_alt"])


def rag_color(sev: str) -> RGBColor:
    m = {"Critical": RED, "High": AMBER, "Medium": AMBER, "Low": GREEN}
    return m.get(sev, MUTED)


def rag_hex(sev: str) -> str:
    m = {"Critical": D.BRAND["red"], "High": D.BRAND["amber"],
         "Medium": D.BRAND["amber"], "Low": D.BRAND["green"]}
    return m.get(sev, D.BRAND["muted"])


# ---------------------------------------------------------------------------
# Low-level drawing helpers
# ---------------------------------------------------------------------------

def _add_rect(slide, x, y, w, h, fill_hex, line_hex=None, line_w=0):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    if line_hex:
        shape.line.color.rgb = rgb(line_hex)
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape


def _add_text(slide, text, x, y, w, h, *,
              size=12, bold=False, italic=False, color=None,
              align=PP_ALIGN.LEFT, valign=None, wrap=True,
              font=None):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    rp = run.font
    rp.size = Pt(size)
    rp.bold = bold
    rp.italic = italic
    rp.color.rgb = color or SLATE
    if font:
        rp.name = font
    return txBox


def _add_paragraph(tf, text, *, size=11, bold=False, italic=False,
                   color=None, align=PP_ALIGN.LEFT, space_before=0, font=None):
    """Append a paragraph to an existing TextFrame."""
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or SLATE
    if font:
        run.font.name = font
    return p


# ---------------------------------------------------------------------------
# Slide chrome: title bar + footer
# ---------------------------------------------------------------------------

TITLE_BAR_H = 0.72   # inches — top accent bar
FOOTER_Y    = 7.05   # inches
FOOTER_H    = 0.38   # inches
CONTENT_TOP = TITLE_BAR_H + 0.18
CONTENT_BOT = FOOTER_Y - 0.12
MARGIN_L    = 0.45
MARGIN_R    = 12.85  # right edge
CONTENT_W   = MARGIN_R - MARGIN_L


def _add_chrome(slide, title_text, slide_num=None, dark_cover=False):
    """Draw the warm page background, top title bar, footer line and slide number."""
    # ── Warm off-white full-bleed page background (Canva-esque) ────────────
    _add_rect(slide, 0, 0, 13.333, 7.5, D.BRAND["band_bg"])

    bar_color = D.BRAND["ink"] if not dark_cover else D.BRAND["ink"]

    # ── Top title bar (deep navy) ──────────────────────────────────────────
    bar = _add_rect(slide, 0, 0, 13.333, TITLE_BAR_H, bar_color)

    # Thin left accent stripe (teal accent)
    stripe_color = D.BRAND["accent"]
    _add_rect(slide, 0, 0, 0.22, TITLE_BAR_H, stripe_color)

    # Title text inside bar
    tx = slide.shapes.add_textbox(
        Inches(0.35), Inches(0.10),
        Inches(11.8), Inches(0.52)
    )
    tx.word_wrap = False
    tf = tx.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = D.BRAND["font_head"]

    # ── Footer ─────────────────────────────────────────────────────────────
    _add_rect(slide, 0, FOOTER_Y, 13.333, FOOTER_H, D.BRAND["band_bg"])
    # thin top rule on footer
    _add_rect(slide, 0, FOOTER_Y, 13.333, 0.012, D.BRAND["rule"])

    footer_tx = slide.shapes.add_textbox(
        Inches(MARGIN_L), Inches(FOOTER_Y + 0.06),
        Inches(10.5), Inches(0.28)
    )
    tf2 = footer_tx.text_frame
    tf2.word_wrap = False
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = f"ESGIntel  |  CONFIDENTIAL  ·  {D.COMPANY['name']}  |  {D.COMPANY.get('as_of', 'ESG Due Diligence')}"
    r2.font.size = Pt(7.5)
    r2.font.color.rgb = MUTED
    r2.font.name = D.BRAND["font_body"]

    if slide_num is not None:
        pg_tx = slide.shapes.add_textbox(
            Inches(12.5), Inches(FOOTER_Y + 0.06),
            Inches(0.7), Inches(0.28)
        )
        tf3 = pg_tx.text_frame
        p3 = tf3.paragraphs[0]
        p3.alignment = PP_ALIGN.RIGHT
        r3 = p3.add_run()
        r3.text = str(slide_num)
        r3.font.size = Pt(7.5)
        r3.font.color.rgb = MUTED
        r3.font.name = D.BRAND["font_body"]


def _action_headline(slide, text):
    """Action-oriented sub-headline just below the title bar."""
    tx = slide.shapes.add_textbox(
        Inches(MARGIN_L), Inches(TITLE_BAR_H + 0.12),
        Inches(CONTENT_W), Inches(0.48)
    )
    tx.word_wrap = True
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(13.5)
    run.font.bold = True
    run.font.color.rgb = INK
    run.font.name = D.BRAND["font_head"]
    return tx


# ---------------------------------------------------------------------------
# Chart helpers (matplotlib → temp PNG → embedded)
# ---------------------------------------------------------------------------

def _save_fig(fig, tmpdir, name) -> str:
    path = os.path.join(tmpdir, name)
    fig.savefig(path, dpi=150, bbox_inches="tight",
                facecolor="white", edgecolor="none")
    plt.close(fig)
    return path


def _chart_pillar_scores(tmpdir) -> str:
    """Horizontal bar chart of E/S/G pillar risk scores."""
    pillars = ["Environmental", "Social", "Governance"]
    scores  = [D.SCORES["pillars"][p]["score"] for p in pillars]
    colors  = []
    for s in scores:
        sev = D.severity(s)
        colors.append(D.hex_rgb01(rag_hex(sev)))

    fig, ax = plt.subplots(figsize=(5.5, 2.4))
    bars = ax.barh(pillars, scores, color=colors, height=0.52, zorder=2)
    ax.set_xlim(0, 100)
    ax.axvline(75, color=D.hex_rgb01(D.BRAND["red"]),   ls="--", lw=1, alpha=0.6)
    ax.axvline(55, color=D.hex_rgb01(D.BRAND["amber"]), ls="--", lw=1, alpha=0.6)
    ax.set_xlabel("Risk Score (0–100, higher = greater risk)", fontsize=7,
                  color=D.hex_rgb01(D.BRAND["muted"])[0:3])
    ax.tick_params(axis="y", labelsize=9)
    ax.tick_params(axis="x", labelsize=7, colors=D.hex_rgb01(D.BRAND["muted"]))
    ax.set_facecolor("white")
    fig.patch.set_facecolor("white")
    ax.grid(axis="x", color="#EAE4D9", zorder=1)
    ax.spines[["top","right","left"]].set_visible(False)
    for bar, score in zip(bars, scores):
        ax.text(bar.get_width() + 1.2, bar.get_y() + bar.get_height()/2,
                str(score), va="center", ha="left", fontsize=9, fontweight="bold",
                color=D.hex_rgb01(D.BRAND["ink"]))
    fig.tight_layout()
    return _save_fig(fig, tmpdir, "pillar_scores.png")


def _chart_financial_exposure(tmpdir) -> str:
    """Horizontal bar chart of financial exposure items."""
    items = D.FINANCIAL_EXPOSURE
    labels = [i["driver"] for i in items]
    amounts = [i["amount_eur_m"] for i in items]
    colors = [D.hex_rgb01(rag_hex(i["rating"])) for i in items]

    # Wrap long labels
    labels = [textwrap.fill(l, 28) for l in labels]

    fig, ax = plt.subplots(figsize=(6.5, 3.2))
    bars = ax.barh(labels, amounts, color=colors, height=0.55, zorder=2)
    ax.set_xlabel("EUR Million", fontsize=7, color=D.hex_rgb01(D.BRAND["muted"]))
    ax.tick_params(axis="y", labelsize=7.5)
    ax.tick_params(axis="x", labelsize=7)
    ax.set_facecolor("white")
    fig.patch.set_facecolor("white")
    ax.grid(axis="x", color="#EAE4D9", zorder=1)
    ax.spines[["top","right","left"]].set_visible(False)
    for bar, amt in zip(bars, amounts):
        ax.text(bar.get_width() + 3, bar.get_y() + bar.get_height()/2,
                f"€{amt}M", va="center", ha="left", fontsize=8, fontweight="bold",
                color=D.hex_rgb01(D.BRAND["ink"]))
    ax.invert_yaxis()
    fig.tight_layout()
    return _save_fig(fig, tmpdir, "financial_exposure.png")


def _chart_climate_scenarios(tmpdir) -> str:
    """Grouped bar / waterfall for climate scenario P&L impacts."""
    scenarios = [s["scenario"] for s in D.CLIMATE_SCENARIOS]
    pnl       = [s["pnl_2030_eur_m"] for s in D.CLIMATE_SCENARIOS]
    carbon    = [s["carbon_2030"] for s in D.CLIMATE_SCENARIOS]

    x = np.arange(len(scenarios))
    colors = [D.hex_rgb01(D.BRAND["red"]),
              D.hex_rgb01(D.BRAND["red"]),
              D.hex_rgb01(D.BRAND["amber"])]

    fig, ax = plt.subplots(figsize=(6.0, 2.8))
    bars = ax.bar(x, pnl, color=colors, width=0.5, zorder=2)
    ax.set_xticks(x)
    ax.set_xticklabels([textwrap.fill(s, 16) for s in scenarios], fontsize=7.5)
    ax.set_ylabel("P&L Impact 2030 (€M)", fontsize=7.5,
                  color=D.hex_rgb01(D.BRAND["muted"]))
    ax.axhline(0, color=D.hex_rgb01(D.BRAND["rule"]), lw=1)
    ax.set_facecolor("white")
    fig.patch.set_facecolor("white")
    ax.grid(axis="y", color="#EAE4D9", zorder=1)
    ax.spines[["top","right","left"]].set_visible(False)
    for bar, val, cp in zip(bars, pnl, carbon):
        ax.text(bar.get_x() + bar.get_width()/2, val - 2,
                f"€{abs(val)}M\n(ETS {cp})",
                ha="center", va="top", fontsize=7, fontweight="bold",
                color="white")
    fig.tight_layout()
    return _save_fig(fig, tmpdir, "climate_scenarios.png")


def _chart_risk_distribution(tmpdir) -> str:
    """Bubble / scatter plot of risks by score, coloured by severity."""
    risks = D.RISKS
    fig, ax = plt.subplots(figsize=(5.5, 2.8))

    pillar_map = {"E": 0, "G": 1, "S": 2}
    pillar_labels = ["Environmental", "Governance", "Social"]
    cmap = {"Critical": D.hex_rgb01(D.BRAND["red"]),
            "High":     D.hex_rgb01(D.BRAND["amber"]),
            "Medium":   D.hex_rgb01(D.BRAND["amber"]),
            "Low":      D.hex_rgb01(D.BRAND["green"])}

    for r in risks:
        y = pillar_map.get(r["pillar"], 0)
        sev = D.severity(r["score"])
        ax.scatter(r["score"], y, s=r["score"]*3, color=cmap[sev],
                   alpha=0.8, zorder=3, edgecolors="white", linewidth=0.6)
        ax.text(r["score"], y + 0.18, r["id"], ha="center", fontsize=6.5,
                color=D.hex_rgb01(D.BRAND["ink"]))

    ax.set_yticks([0, 1, 2])
    ax.set_yticklabels(pillar_labels, fontsize=8)
    ax.set_xlabel("Risk Score →  (higher = greater risk)", fontsize=7.5,
                  color=D.hex_rgb01(D.BRAND["muted"]))
    ax.set_xlim(20, 100)
    ax.set_ylim(-0.5, 2.7)
    ax.axvline(75, color=D.hex_rgb01(D.BRAND["red"]),   ls="--", lw=0.9, alpha=0.5)
    ax.axvline(55, color=D.hex_rgb01(D.BRAND["amber"]), ls="--", lw=0.9, alpha=0.5)
    ax.set_facecolor("white")
    fig.patch.set_facecolor("white")
    ax.spines[["top","right"]].set_visible(False)
    ax.grid(axis="x", color="#EAE4D9", zorder=1)
    # Legend patches
    patches = [
        mpatches.Patch(color=D.hex_rgb01(D.BRAND["red"]),   label="Critical (≥75)"),
        mpatches.Patch(color=D.hex_rgb01(D.BRAND["amber"]), label="High/Med (35–74)"),
    ]
    ax.legend(handles=patches, fontsize=6.5, loc="lower right", framealpha=0.7)
    fig.tight_layout()
    return _save_fig(fig, tmpdir, "risk_distribution.png")


def _chart_physical_vs_transition(tmpdir) -> str:
    """Side-by-side bars: physical risks vs transition risks."""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(7.5, 2.8))

    # Physical
    ph = sorted(D.PHYSICAL_RISKS, key=lambda x: x["score"], reverse=True)
    ph_labels = [r["hazard"] for r in ph]
    ph_scores = [r["score"] for r in ph]
    ph_colors = [D.hex_rgb01(rag_hex(D.severity(s))) for s in ph_scores]
    ax1.barh(ph_labels, ph_scores, color=ph_colors, height=0.5, zorder=2)
    ax1.set_xlim(0, 100)
    ax1.set_title("Physical Risks", fontsize=9, fontweight="bold",
                  color=D.hex_rgb01(D.BRAND["ink"]))
    ax1.tick_params(axis="y", labelsize=7)
    ax1.tick_params(axis="x", labelsize=6.5)
    ax1.set_facecolor("white")
    ax1.spines[["top","right","left"]].set_visible(False)
    ax1.grid(axis="x", color="#EAE4D9", zorder=1)
    ax1.invert_yaxis()

    # Transition
    tr = sorted(D.TRANSITION_RISKS, key=lambda x: x["score"], reverse=True)
    tr_labels = [r["driver"] for r in tr]
    tr_scores = [r["score"] for r in tr]
    tr_colors = [D.hex_rgb01(rag_hex(D.severity(s))) for s in tr_scores]
    ax2.barh(tr_labels, tr_scores, color=tr_colors, height=0.5, zorder=2)
    ax2.set_xlim(0, 100)
    ax2.set_title("Transition Risks", fontsize=9, fontweight="bold",
                  color=D.hex_rgb01(D.BRAND["ink"]))
    ax2.tick_params(axis="y", labelsize=7)
    ax2.tick_params(axis="x", labelsize=6.5)
    ax2.set_facecolor("white")
    ax2.spines[["top","right","left"]].set_visible(False)
    ax2.grid(axis="x", color="#EAE4D9", zorder=1)
    ax2.invert_yaxis()

    fig.patch.set_facecolor("white")
    fig.tight_layout()
    return _save_fig(fig, tmpdir, "physical_transition.png")


# ---------------------------------------------------------------------------
# Reusable slide-content builders
# ---------------------------------------------------------------------------

def _big_number_tile(slide, value, label, x, y, w=2.5, h=1.3,
                     bg_hex=None, val_color=None):
    """A 'KPI tile' with a large number and a small label below it."""
    bg_hex = bg_hex or D.BRAND["accent_lt"]
    val_color = val_color or D.BRAND["ink"]
    _add_rect(slide, x, y, w, h, bg_hex)
    # value
    _add_text(slide, value,
              x + 0.12, y + 0.08, w - 0.24, h * 0.55,
              size=22, bold=True, color=rgb(val_color),
              align=PP_ALIGN.CENTER, font=D.BRAND["font_head"])
    # label
    _add_text(slide, label,
              x + 0.12, y + h * 0.57, w - 0.24, h * 0.38,
              size=8, bold=False, color=MUTED,
              align=PP_ALIGN.CENTER, font=D.BRAND["font_body"])


def _rag_chip(slide, text, severity, x, y, w=1.7, h=0.34):
    """A coloured RAG chip / badge."""
    h_color = rag_hex(severity)
    _add_rect(slide, x, y, w, h, h_color)
    _add_text(slide, text,
              x + 0.08, y + 0.03, w - 0.16, h - 0.06,
              size=8, bold=True, color=WHITE,
              align=PP_ALIGN.CENTER, font=D.BRAND["font_body"])


def _section_header(slide, label):
    """Teal pill + label — visual section divider within a slide."""
    _add_rect(slide, MARGIN_L, TITLE_BAR_H + 0.62, 0.22, 0.24, ACCENT)
    _add_text(slide, label,
              MARGIN_L + 0.30, TITLE_BAR_H + 0.62, CONTENT_W, 0.28,
              size=9, bold=True, color=ACCENT,
              font=D.BRAND["font_body"])


def _section_divider_slide(prs, section_num: int, title: str, subtitle: str = ""):
    """
    Full-bleed dark navy section-divider slide (Canva reference style).
    Large teal section number on left, section title + subtitle in white.
    """
    blank_layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank_layout)

    # Full dark-navy background
    _add_rect(s, 0, 0, 13.333, 7.5, D.BRAND["ink"])

    # Teal left accent stripe
    _add_rect(s, 0, 0, 0.35, 7.5, D.BRAND["accent"])

    # Subtle warm-cream horizontal rule at 50%
    _add_rect(s, 0.55, 3.6, 12.5, 0.025, D.BRAND.get("rule", "CCCCCC"))

    # Large section number (e.g. "01") in teal
    num_str = f"{section_num:02d}"
    _add_text(s, num_str, 0.65, 1.2, 3.5, 2.8,
              size=120, bold=True, color=ACCENT,
              font=D.BRAND["font_head"])

    # Section title — large white text
    _add_text(s, title, 4.2, 2.05, 8.6, 1.4,
              size=36, bold=True, color=WHITE,
              font=D.BRAND["font_head"])

    # Optional subtitle — smaller muted text
    if subtitle:
        _add_text(s, subtitle, 4.2, 3.55, 8.6, 0.75,
                  size=14, bold=False, color=MUTED,
                  font=D.BRAND["font_body"])

    # Footer branding
    _add_text(s, "ESGIntel  |  CONFIDENTIAL", 0.65, 6.9, 12.0, 0.35,
              size=8, bold=False, color=MUTED,
              font=D.BRAND["font_body"])

    return s


def _add_table(slide, headers, rows, x, y, w, h,
               col_widths=None, font_size=8.5):
    """Add a styled table with shaded header row."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor as RC

    n_cols = len(headers)
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(n_rows, n_cols,
                                  Inches(x), Inches(y),
                                  Inches(w), Inches(h)).table

    # Set column widths
    if col_widths:
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = Inches(cw)

    def _set_cell(cell, text, *, bold=False, bg=None, fg=None, align=PP_ALIGN.LEFT):
        cell.text = text
        tf = cell.text_frame
        tf.word_wrap = True
        for para in tf.paragraphs:
            para.alignment = align
            for run in para.runs:
                run.font.size = Pt(font_size)
                run.font.bold = bold
                run.font.color.rgb = fg or SLATE
                run.font.name = D.BRAND["font_body"]
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

    # Header row
    for ci, hdr in enumerate(headers):
        _set_cell(tbl.cell(0, ci), hdr,
                  bold=True, bg=ACCENT, fg=WHITE, align=PP_ALIGN.CENTER)

    # Data rows (alternating warm-gray)
    for ri, row in enumerate(rows):
        bg = ROW_ALT if ri % 2 == 1 else PAGE_BG
        for ci, val in enumerate(row):
            _set_cell(tbl.cell(ri + 1, ci), str(val), bg=bg)

    return tbl


# ---------------------------------------------------------------------------
# ═══════════════════════════════════════════════════════════════
# BOARD DECK  (~15 slides)
# ═══════════════════════════════════════════════════════════════
# ---------------------------------------------------------------------------

def _build_board_deck(tmpdir: str, out_path: str):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # fully blank

    today = datetime.date.today().strftime("%B %Y")

    # ── Slide 1: Cover (warm off-white, teal accent bar) ──────────────────
    s = prs.slides.add_slide(blank_layout)
    # Warm cream full-bleed background
    _add_rect(s, 0, 0, 13.333, 7.5, D.BRAND["cream"])
    # Bold teal accent bar on left
    _add_rect(s, 0, 0, 0.55, 7.5, D.BRAND["accent"])
    # Teal underline beneath the title block
    _add_rect(s, 0.85, 2.95, 5.6, 0.06, D.BRAND["accent"])

    _add_text(s, "ESGIntel", 0.85, 0.65, 11.0, 0.6,
              size=12, bold=True, color=ACCENT, font=D.BRAND["font_body"])
    _add_text(s, D.COMPANY["name"], 0.85, 1.25, 11.5, 1.3,
              size=44, bold=True, color=rgb(D.BRAND["ink"]), font=D.BRAND["font_head"])
    _add_text(s, "ESG Engagement Deck — Board Briefing", 0.85, 3.15, 11.0, 0.7,
              size=20, bold=False, color=ACCENT, font=D.BRAND["font_head"])
    _add_text(s, D.BRAND["confidential"], 0.85, 3.95, 11.0, 0.45,
              size=9, bold=False, color=rgb(D.BRAND["muted"]),
              font=D.BRAND["font_body"])
    _add_text(s, today, 0.85, 4.42, 3.5, 0.38,
              size=9, color=rgb(D.BRAND["muted"]), font=D.BRAND["font_body"])

    # Bottom metadata
    meta_y = 5.35
    for label, val in [
        ("Sector",    D.COMPANY["sector"]),
        ("Domicile",  D.COMPANY["domicile"]),
        ("Employees", f"{D.COMPANY['employees']:,}"),
        ("Revenue",   f"€{D.COMPANY['revenue_eur_bn']}bn ({D.COMPANY['report_period']})"),
    ]:
        _add_text(s, label.upper(), 0.85, meta_y, 2.0, 0.25,
                  size=7, bold=True, color=ACCENT, font=D.BRAND["font_body"])
        _add_text(s, val, 2.75, meta_y, 9.5, 0.28,
                  size=8.5, color=rgb(D.BRAND["slate"]), font=D.BRAND["font_body"])
        meta_y += 0.38

    _add_text(s, "CONFIDENTIAL", 9.5, 7.1, 3.5, 0.35,
              size=7.5, bold=True, color=rgb(D.BRAND["muted"]),
              align=PP_ALIGN.RIGHT, font=D.BRAND["font_body"])

    # ── Section 01 Divider: ESG Overview ──────────────────────────────────
    _section_divider_slide(prs, 1, "ESG Overview",
                           "Executive Summary  ·  Scorecard  ·  Risk Register")

    # ── Slide 2: Executive Summary ────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    _add_chrome(s, "Executive Summary", 2)
    _action_headline(s, f"{D.COMPANY['short']} faces a compound ESG risk profile with €950M+ in quantifiable near-term exposure")

    # Three KPI tiles
    tile_y = TITLE_BAR_H + 0.72
    _big_number_tile(s, "38 / 100", "ESG Risk Score", MARGIN_L, tile_y,
                     w=2.8, h=1.45, bg_hex=D.BRAND["red"], val_color="FFFFFF")
    _big_number_tile(s, "P19", "vs NACE C24 Peers", MARGIN_L + 2.95, tile_y,
                     w=2.8, h=1.45, bg_hex=D.BRAND["amber"], val_color="FFFFFF")
    _big_number_tile(s, "€950M+", "Quantified Exposure", MARGIN_L + 5.90, tile_y,
                     w=2.8, h=1.45, bg_hex=D.BRAND["accent"], val_color="FFFFFF")
    _big_number_tile(s, "72%", "Data Confidence", MARGIN_L + 8.85, tile_y,
                     w=2.8, h=1.45, bg_hex=D.BRAND["accent_lt"], val_color=D.BRAND["ink"])

    # Key messages
    msg_y = tile_y + 1.65
    messages = [
        ("Carbon cost is the single biggest risk", "EU ETS at €68M/yr net by 2030; CBAM adds €24M — no hedging policy exists."),
        ("CSRD first report due FY2025 — readiness is low", "Double-materiality assessment incomplete; ESRS E1 transition plan absent."),
        ("Three governance gaps demand immediate action", "No whistleblower mechanism, no ESG pay link, no SBTi commitment."),
        ("Green finance upside of €450–600M available", "EU Green Bond + SLL + Innovation Fund — contingent on ESG remediation."),
    ]
    for i, (hdr, body) in enumerate(messages):
        mx = MARGIN_L if i < 2 else MARGIN_L + 6.3
        my = msg_y if i % 2 == 0 else msg_y + 1.25
        _add_rect(s, mx, my, 5.85, 1.1, D.BRAND["band_bg"])
        _add_rect(s, mx, my, 0.06, 1.1, D.BRAND["accent"])
        _add_text(s, hdr, mx + 0.16, my + 0.06, 5.5, 0.32,
                  size=9.5, bold=True, color=INK, font=D.BRAND["font_head"])
        _add_text(s, body, mx + 0.16, my + 0.42, 5.5, 0.6,
                  size=8.5, color=SLATE, font=D.BRAND["font_body"])

    # ── Slide 3: ESG Scorecard ────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    _add_chrome(s, "ESG Scorecard", 3)
    _action_headline(s, "Bottom-quartile risk profile across all three pillars — Environmental exposure dominates")

    # Overall score prominent
    _big_number_tile(s, "38 / 100", "Overall ESG Risk Score", MARGIN_L, TITLE_BAR_H + 0.70,
                     w=3.0, h=1.55, bg_hex=D.BRAND["red"], val_color="FFFFFF")
    _add_text(s, "High Risk — P19 vs NACE C24 Peers (n=47)",
              MARGIN_L + 3.15, TITLE_BAR_H + 0.70, 5.8, 0.45,
              size=12, bold=True, color=RED, font=D.BRAND["font_head"])
    _add_text(s, f"Score interpretation: higher = greater ESG risk\n"
                 f"Formula: (E×40%) + (S×30%) + (G×30%)\n"
                 f"Data confidence: {D.SCORES['confidence']}%  |  Reference period: {D.COMPANY['report_period']}",
              MARGIN_L + 3.15, TITLE_BAR_H + 1.22, 5.8, 0.95,
              size=8.5, color=MUTED, font=D.BRAND["font_body"])

    # Pillar chart
    chart_path = _chart_pillar_scores(tmpdir)
    s.shapes.add_picture(chart_path, Inches(MARGIN_L + 3.15), Inches(TITLE_BAR_H + 2.20),
                         Inches(5.8), Inches(2.75))

    # Pillar detail tiles
    tile_x_start = MARGIN_L
    tile_y2 = TITLE_BAR_H + 2.40
    for pillar, data in D.SCORES["pillars"].items():
        sev = data["label"]
        _add_rect(s, tile_x_start, tile_y2, 2.85, 2.45, D.BRAND["band_bg"])
        _add_rect(s, tile_x_start, tile_y2, 2.85, 0.06, rag_hex(sev))
        _add_text(s, pillar, tile_x_start + 0.12, tile_y2 + 0.12, 2.6, 0.38,
                  size=10, bold=True, color=INK, font=D.BRAND["font_head"])
        _add_text(s, f"{data['score']} / 100",
                  tile_x_start + 0.12, tile_y2 + 0.55, 2.6, 0.55,
                  size=20, bold=True, color=rgb(rag_hex(sev)),
                  font=D.BRAND["font_head"])
        _add_text(s, f"{sev}  |  P{data['percentile']} vs peers",
                  tile_x_start + 0.12, tile_y2 + 1.15, 2.6, 0.32,
                  size=8, color=MUTED, font=D.BRAND["font_body"])
        _add_text(s, f"Weight: {int(data['weight']*100)}%",
                  tile_x_start + 0.12, tile_y2 + 1.52, 2.6, 0.25,
                  size=7.5, color=MUTED, font=D.BRAND["font_body"])
        tile_x_start += 3.0

    # ── Slide 4: Double-Materiality ───────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    _add_chrome(s, "Double-Materiality Framing", 4)
    _action_headline(s, f"CSRD requires {D.COMPANY['short']} to assess both financial impact and societal impact of ESG topics")

    # Two-column explainer
    col_y = TITLE_BAR_H + 0.72
    for col, (title, color_h, body_items) in enumerate([
        ("Financial Materiality (Inside-Out)", D.BRAND["accent"],
         [f"How ESG risks and opportunities affect {D.COMPANY['short']}'s cash flows, assets and cost of capital",
          "ETS carbon costs: €68M/yr net by 2030",
          "CBAM certificates: €24M/yr by 2027",
          "Asset stranding: €275M impairment risk",
          "DRI/EAF capex: €500M phased transition",
          "Green finance upside: €450–600M"]),
        ("Impact Materiality (Outside-In)", D.BRAND["red"],
         [f"How {D.COMPANY['short']}'s operations impact people, communities and the environment",
          "842,000 tCO2e Scope 1 emissions — steel sector's contribution to Paris 1.5°C",
          "Water withdrawal at high-stress Bilbao & Taranto basins",
          "High-risk supply chain (Brazil iron ore, DRC minerals)",
          "LTIFR 33% above sector median; 1 fatality in FY2022",
          "Whistleblower gap: 8,400 employees without a reporting channel"]),
    ]):
        cx = MARGIN_L + col * 6.3
        _add_rect(s, cx, col_y, 6.0, 4.65, D.BRAND["band_bg"])
        _add_rect(s, cx, col_y, 6.0, 0.06, color_h)
        _add_text(s, title, cx + 0.15, col_y + 0.12, 5.7, 0.45,
                  size=11, bold=True, color=INK, font=D.BRAND["font_head"])
        for bi, item in enumerate(body_items):
            iy = col_y + 0.65 + bi * 0.60
            prefix = "→  " if bi > 0 else ""
            _add_text(s, prefix + item, cx + 0.15, iy, 5.7, 0.56,
                      size=8.5, bold=(bi == 0), color=SLATE if bi > 0 else INK,
                      font=D.BRAND["font_body"])

    # DMA status note
    note_y = col_y + 4.80
    _add_rect(s, MARGIN_L, note_y, CONTENT_W, 0.50, "FFF3CD")
    _add_text(s, f"DMA Status: {D.COMPANY['short']}'s double-materiality assessment is INCOMPLETE — critical gap for CSRD FY2025 first report.",
              MARGIN_L + 0.15, note_y + 0.08, CONTENT_W - 0.3, 0.35,
              size=8.5, bold=True, color=rgb(D.BRAND["amber"]),
              font=D.BRAND["font_body"])

    # ── Slide 5: Top Risks Heatmap ────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    _add_chrome(s, "Top Risk Register", 5)
    _action_headline(s, "Three critical-rated risks demand immediate Board attention — all with direct P&L exposure")

    chart_path = _chart_risk_distribution(tmpdir)
    s.shapes.add_picture(chart_path, Inches(MARGIN_L), Inches(TITLE_BAR_H + 0.65),
                         Inches(5.2), Inches(2.65))

    # Risk table (right side)
    headers = ["ID", "Title", "Score", "Severity", "Horizon", "Financial"]
    rows = []
    for r in sorted(D.RISKS, key=lambda x: x["score"], reverse=True)[:8]:
        rows.append([
            r["id"],
            r["title"][:45] + ("…" if len(r["title"]) > 45 else ""),
            str(r["score"]),
            D.severity(r["score"]),
            r["horizon"][:18],
            r["financial"][:22],
        ])
    _add_table(s, headers, rows,
               MARGIN_L + 5.4, TITLE_BAR_H + 0.65,
               CONTENT_W - 5.4, 5.85,
               col_widths=[0.6, 2.6, 0.55, 0.72, 1.25, 1.45],
               font_size=7.5)

    # ── Section 02 Divider: Pillar Deep-Dives ─────────────────────────────
    _section_divider_slide(prs, 2, "Pillar Deep-Dives",
                           "Environmental  ·  Social  ·  Governance")

    # ── Slide 6: Environmental Deep-Dive ──────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    _add_chrome(s, "Environmental — E Pillar Deep-Dive", 6)
    _action_headline(s, "Score 64/100 (High, P81) — Carbon intensity and water stress are structural exposures")

    env_kpis = [k for k in D.KPIS if k["group"].startswith("Climate") or k["group"].startswith("Water") or k["group"].startswith("Circular")]
    _add_table(s, ["Metric", f"{D.COMPANY['short']}", "Peer", "Pctl", "Direction"],
               [[k["metric"], k["value"], str(k["peer"]), str(k["pctl"]) if k["pctl"] else "N/A",
                 "▼ Lower better" if k["dir"] == "lower_better" else "▲ Higher better"]
                for k in env_kpis],
               MARGIN_L, TITLE_BAR_H + 0.72, 6.5, 3.2,
               col_widths=[2.4, 1.6, 1.1, 0.65, 1.1], font_size=8.5)

    # Key insights
    insights = [
        ("Scope 1: 842 ktCO2e", "P89 vs peers — highest decile of carbon intensity. BF-BOF process locks in ~1.8-2.1 tCO2e/t."),
        ("Water recycling: 12%", "vs 45-60% best practice (P12). Bilbao permit review 2027. €18M revenue at risk at 30% cut."),
        ("EAF scrap: 71%", "Above sector median (62%). One relative bright spot; scrap availability limits further upside."),
        ("Renewable electricity", "Not disclosed — blocks green-bond eligibility and Scope 2 reduction narrative."),
    ]
    insight_y = TITLE_BAR_H + 4.10
    for i, (hdr, body) in enumerate(insights):
        ix = MARGIN_L + (i % 2) * 6.3
        iy = insight_y + (i // 2) * 1.05
        _add_rect(s, ix, iy, 5.95, 0.95, D.BRAND["accent_lt"])
        _add_rect(s, ix, iy, 0.05, 0.95, D.BRAND["accent"])
        _add_text(s, hdr, ix + 0.14, iy + 0.06, 5.65, 0.30,
                  size=9, bold=True, color=INK, font=D.BRAND["font_head"])
        _add_text(s, body, ix + 0.14, iy + 0.40, 5.65, 0.50,
                  size=8, color=SLATE, font=D.BRAND["font_body"])

    # ── Slide 7: Social Deep-Dive ─────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    _add_chrome(s, "Social — S Pillar Deep-Dive", 7)
    _action_headline(s, "Score 57/100 (Medium, P63) — H&S performance and HRDD absence are the critical gaps")

    soc_kpis = [k for k in D.KPIS if k["group"].startswith("Workforce")]
    _add_table(s, ["Metric", f"{D.COMPANY['short']}", "Peer", "Pctl", "Direction"],
               [[k["metric"], k["value"], str(k["peer"]), str(k["pctl"]) if k["pctl"] else "N/A",
                 "▼ Lower better" if k["dir"] == "lower_better" else "▲ Higher better"]
                for k in soc_kpis],
               MARGIN_L, TITLE_BAR_H + 0.72, 6.5, 3.2,
               col_widths=[2.4, 1.6, 1.1, 0.65, 1.1], font_size=8.5)

    soc_risks = [r for r in D.RISKS if r["pillar"] == "S"]
    risk_y = TITLE_BAR_H + 4.10
    for i, r in enumerate(soc_risks[:4]):
        rx = MARGIN_L + (i % 2) * 6.3
        ry = risk_y + (i // 2) * 1.05
        sev = D.severity(r["score"])
        _add_rect(s, rx, ry, 5.95, 0.95, D.BRAND["band_bg"])
        _rag_chip(s, sev, sev, rx + 0.12, ry + 0.08, w=1.35, h=0.28)
        _add_text(s, r["title"], rx + 1.6, ry + 0.06, 4.2, 0.32,
                  size=8.5, bold=True, color=INK, font=D.BRAND["font_head"])
        _add_text(s, r["summary"][:120] + "…", rx + 0.12, ry + 0.45, 5.7, 0.45,
                  size=7.5, color=SLATE, font=D.BRAND["font_body"])

    # ── Slide 8: Governance Deep-Dive ─────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    _add_chrome(s, "Governance — G Pillar Deep-Dive", 8)
    _action_headline(s, "Score 49/100 (Medium, P58) — Two immediate compliance failures; ESG pay link absent")

    gov_kpis = [k for k in D.KPIS if k["group"].startswith("Governance")]
    _add_table(s, ["Metric", f"{D.COMPANY['short']}", "Peer", "Pctl", "Status"],
               [[k["metric"], k["value"], str(k["peer"]), str(k["pctl"]) if k["pctl"] else "N/A",
                 "OK" if (k["pctl"] or 0) >= 50 else "GAP"]
                for k in gov_kpis],
               MARGIN_L, TITLE_BAR_H + 0.72, 6.5, 2.0,
               col_widths=[2.4, 1.6, 1.4, 0.65, 0.7], font_size=8.5)

    gov_risks = [r for r in D.RISKS if r["pillar"] == "G"]
    risk_y = TITLE_BAR_H + 3.00
    for i, r in enumerate(gov_risks):
        rx = MARGIN_L + (i % 2) * 6.3
        ry = risk_y + (i // 2) * 1.45
        sev = D.severity(r["score"])
        _add_rect(s, rx, ry, 5.95, 1.35, D.BRAND["band_bg"])
        _add_rect(s, rx, ry, 5.95, 0.05, rag_hex(sev))
        _rag_chip(s, f"{sev} ({r['score']})", sev, rx + 0.12, ry + 0.12, w=1.6, h=0.28)
        _add_text(s, r["title"], rx + 1.85, ry + 0.10, 3.95, 0.32,
                  size=9, bold=True, color=INK, font=D.BRAND["font_head"])
        _add_text(s, r["summary"][:160] + "…", rx + 0.12, ry + 0.52, 5.7, 0.55,
                  size=7.5, color=SLATE, font=D.BRAND["font_body"])
        _add_text(s, "Action: " + r["action"][:80], rx + 0.12, ry + 1.08, 5.7, 0.22,
                  size=7, bold=False, italic=True, color=ACCENT,
                  font=D.BRAND["font_body"])

    # ── Section 03 Divider: Climate Risk ──────────────────────────────────
    _section_divider_slide(prs, 3, "Climate Risk",
                           "Physical risks  ·  Transition risks  ·  NGFS Scenarios")

    # ── Slide 9: Climate Physical vs Transition ────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    _add_chrome(s, "Climate Risk — Physical & Transition", 9)
    _action_headline(s, "Transition risks dominate near-term (ETS/CBAM); physical risks compound from 2030 via water & flood")

    chart_path = _chart_physical_vs_transition(tmpdir)
    s.shapes.add_picture(chart_path, Inches(MARGIN_L), Inches(TITLE_BAR_H + 0.68),
                         Inches(CONTENT_W), Inches(3.15))

    # Key physical risk call-outs
    ph_y = TITLE_BAR_H + 4.00
    _add_rect(s, MARGIN_L, ph_y, CONTENT_W, 0.03, D.BRAND["rule"])
    _add_text(s, "KEY PHYSICAL RISK SITES", MARGIN_L, ph_y + 0.08, CONTENT_W, 0.25,
              size=8, bold=True, color=ACCENT, font=D.BRAND["font_body"])
    for i, ph in enumerate(D.PHYSICAL_RISKS[:4]):
        px = MARGIN_L + i * 3.0
        py = ph_y + 0.38
        sev = D.severity(ph["score"])
        _add_rect(s, px, py, 2.8, 0.9, D.BRAND["band_bg"])
        _rag_chip(s, str(ph["score"]), sev, px + 0.1, py + 0.08, w=0.7, h=0.26)
        _add_text(s, ph["hazard"], px + 0.9, py + 0.06, 1.8, 0.28,
                  size=8, bold=True, color=INK, font=D.BRAND["font_body"])
        _add_text(s, ph["sites"], px + 0.1, py + 0.42, 2.6, 0.36,
                  size=7.5, color=MUTED, font=D.BRAND["font_body"])

    # ── Slide 10: Climate Scenario P&L ────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    _add_chrome(s, "Climate Scenario Analysis — P&L Impact 2030", 10)
    _action_headline(s, f"Under all NGFS scenarios {D.COMPANY['short']} faces negative P&L; disorderly transition is most severe at -€96M")

    chart_path = _chart_climate_scenarios(tmpdir)
    s.shapes.add_picture(chart_path, Inches(MARGIN_L), Inches(TITLE_BAR_H + 0.68),
                         Inches(6.2), Inches(3.1))

    # Scenario detail table
    headers = ["Scenario", "ETS Price 2030", "P&L Impact 2030", "Key Driver"]
    rows = []
    for sc in D.CLIMATE_SCENARIOS:
        rows.append([sc["scenario"], sc["carbon_2030"],
                     f"€{sc['pnl_2030_eur_m']}M", sc["note"]])
    _add_table(s, headers, rows,
               MARGIN_L + 6.4, TITLE_BAR_H + 0.68, CONTENT_W - 6.4, 2.2,
               col_widths=[2.3, 1.2, 1.3, 1.8], font_size=8)

    # Methodology note
    note_y = TITLE_BAR_H + 3.95
    _add_rect(s, MARGIN_L, note_y, CONTENT_W, 1.50, D.BRAND["band_bg"])
    _add_text(s, "Scenario Methodology (NGFS-aligned)", MARGIN_L + 0.15, note_y + 0.10,
              CONTENT_W - 0.3, 0.30, size=9, bold=True, color=INK,
              font=D.BRAND["font_head"])
    _add_text(s,
              "Orderly (Net Zero 2050): rapid, smooth transition — high near-term carbon cost but rewarded via capex deployment and market access.\n"
              "Disorderly (Delayed): policy delay followed by sharp repricing — highest stranding risk; BF-BOF assets unviable by 2032.\n"
              "Hot House (NDCs only): low carbon cost, but physical damage (water, flood) compounds from 2030 — Taranto and Bilbao most exposed.\n"
              "All scenarios exclude CBAM exposure (additive) and DRI/EAF transition capex (€400–600M).",
              MARGIN_L + 0.15, note_y + 0.45, CONTENT_W - 0.3, 1.0,
              size=8, color=SLATE, font=D.BRAND["font_body"])

    # ── Slide 11: Financial Exposure Bridge ───────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    _add_chrome(s, "Financial Exposure Bridge", 11)
    _action_headline(s, "Quantified ESG exposures total >€950M — ETS and asset stranding are the largest line items")

    chart_path = _chart_financial_exposure(tmpdir)
    s.shapes.add_picture(chart_path, Inches(MARGIN_L), Inches(TITLE_BAR_H + 0.65),
                         Inches(6.8), Inches(3.55))

    # Summary table
    _add_table(s, ["Driver", "€M", "Type", "Rating"],
               [[f["driver"], f"€{f['amount_eur_m']}M", f["type"], f["rating"]]
                for f in D.FINANCIAL_EXPOSURE],
               MARGIN_L + 7.0, TITLE_BAR_H + 0.65, CONTENT_W - 7.0, 3.55,
               col_widths=[2.5, 0.7, 1.2, 0.8], font_size=7.5)

    # Total callout
    total_y = TITLE_BAR_H + 4.38
    total = sum(f["amount_eur_m"] for f in D.FINANCIAL_EXPOSURE)
    _big_number_tile(s, f"€{total}M", "Total Quantified Exposure (exc. phased capex)",
                     MARGIN_L, total_y, w=6.8, h=0.85,
                     bg_hex=D.BRAND["red"], val_color="FFFFFF")

    # ── Section 04 Divider: Regulatory & Finance ──────────────────────────
    _section_divider_slide(prs, 4, "Regulatory & Green Finance",
                           "CSRD  ·  ETS/CBAM  ·  Green Bond  ·  SLL")

    # ── Slide 12: Regulatory Timeline ─────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    _add_chrome(s, "Regulatory Timeline", 12)
    _action_headline(s, f"{D.COMPANY['short']} has 18 months to its first CSRD report — five overlapping regulatory clocks are ticking")

    tl_y = TITLE_BAR_H + 0.82
    tl_h = 0.22
    tl_line_y = tl_y + tl_h / 2
    # Horizontal timeline spine
    _add_rect(s, MARGIN_L + 0.8, tl_line_y, CONTENT_W - 0.8, 0.04, D.BRAND["rule"])

    years = sorted(set(r["year"] for r in D.REG_TIMELINE))
    x_positions = {yr: MARGIN_L + 0.8 + i * (CONTENT_W - 0.8) / max(len(years) - 1, 1)
                   for i, yr in enumerate(years)}

    # Year markers
    for yr, xp in x_positions.items():
        _add_rect(s, xp - 0.04, tl_line_y - 0.04, 0.08, 0.08 + tl_h/2, D.BRAND["accent"])
        _add_text(s, yr, xp - 0.5, tl_y - 0.32, 1.0, 0.26,
                  size=10, bold=True, color=ACCENT, align=PP_ALIGN.CENTER,
                  font=D.BRAND["font_head"])

    # Event cards (alternating above/below)
    for i, item in enumerate(D.REG_TIMELINE):
        xp = x_positions[item["year"]]
        above = (i % 2 == 0)
        card_y = (tl_y - 1.45) if above else (tl_y + tl_h + 0.28)
        card_w = 2.6
        status_colors = {"Active": D.BRAND["red"], "In scope": D.BRAND["amber"],
                         "Imminent": D.BRAND["red"], "Upcoming": D.BRAND["amber"],
                         "Non-compliant today": D.BRAND["red"]}
        s_color = status_colors.get(item["status"], D.BRAND["accent"])
        _add_rect(s, xp - card_w/2, card_y, card_w, 1.28, D.BRAND["band_bg"])
        _add_rect(s, xp - card_w/2, card_y, card_w, 0.05, s_color)
        _add_text(s, item["status"], xp - card_w/2 + 0.1, card_y + 0.10,
                  card_w - 0.2, 0.24,
                  size=7, bold=True, color=rgb(s_color), font=D.BRAND["font_body"])
        _add_text(s, item["item"], xp - card_w/2 + 0.1, card_y + 0.38,
                  card_w - 0.2, 0.82,
                  size=7.5, color=SLATE, font=D.BRAND["font_body"])

    # ── Slide 13: Green Finance Opportunity ───────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    _add_chrome(s, "Green Finance Opportunity", 13)
    _action_headline(s, "ESG remediation unlocks €450–600M in lower-cost financing — contingent on SBTi and CSRD readiness")

    # Three opportunity cards
    card_w = 3.9
    for i, opp in enumerate(D.OPPORTUNITIES):
        cx = MARGIN_L + i * (card_w + 0.2)
        cy = TITLE_BAR_H + 0.72
        _add_rect(s, cx, cy, card_w, 4.2, D.BRAND["band_bg"])
        _add_rect(s, cx, cy, card_w, 0.08, D.BRAND["accent"])
        _add_text(s, opp["lever"], cx + 0.15, cy + 0.14, card_w - 0.3, 0.52,
                  size=12, bold=True, color=INK, font=D.BRAND["font_head"])
        _big_number_tile(s, opp["size"], "Financing Size",
                         cx + 0.15, cy + 0.75, card_w - 0.3, 0.95,
                         bg_hex=D.BRAND["accent_lt"], val_color=D.BRAND["ink"])
        _add_text(s, "Benefit:", cx + 0.15, cy + 1.82, card_w - 0.3, 0.25,
                  size=8, bold=True, color=INK, font=D.BRAND["font_body"])
        _add_text(s, opp["benefit"], cx + 0.15, cy + 2.10, card_w - 0.3, 0.65,
                  size=8.5, color=SLATE, font=D.BRAND["font_body"])
        _add_rect(s, cx + 0.15, cy + 2.82, card_w - 0.3, 0.03, D.BRAND["rule"])
        _add_text(s, "Pre-condition:", cx + 0.15, cy + 2.90, card_w - 0.3, 0.25,
                  size=8, bold=True, color=AMBER, font=D.BRAND["font_body"])
        _add_text(s, opp["precondition"], cx + 0.15, cy + 3.18, card_w - 0.3, 0.90,
                  size=8, color=SLATE, font=D.BRAND["font_body"])

    # ── Section 05 Divider: Engagement Asks ───────────────────────────────
    _section_divider_slide(prs, 5, "Engagement Asks",
                           "Prioritised actions for Board & Management")

    # ── Slide 14: Engagement Asks ─────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    _add_chrome(s, "Engagement Asks", 14)
    _action_headline(s, "Six prioritised asks — two demand Board action within 60 days")

    headers = ["Priority", "Ask", "Owner", "Timeline"]
    rows = [[ea["priority"], ea["ask"], ea["owner"], ea["by"]]
            for ea in D.ENGAGEMENT_ASKS]
    _add_table(s, headers, rows,
               MARGIN_L, TITLE_BAR_H + 0.72, CONTENT_W, 5.5,
               col_widths=[0.88, 6.0, 2.3, 1.5], font_size=9)

    # ── Slide 15: Closing ─────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    _add_rect(s, 0, 0, 13.333, 7.5, D.BRAND["ink"])
    _add_rect(s, 0, 0, 0.42, 7.5, D.BRAND["accent"])
    _add_rect(s, 0.42, 3.2, 12.91, 0.04, D.BRAND["accent"])

    _add_text(s, "ESGIntel", 0.7, 0.55, 10.0, 0.5,
              size=11, color=ACCENT, font=D.BRAND["font_body"])
    _add_text(s, "Next Steps & Contact", 0.7, 1.15, 10.0, 0.9,
              size=32, bold=True, color=WHITE, font=D.BRAND["font_head"])

    steps = [
        "1.  Board adopts Engagement Asks as formal agenda items — assign owners.",
        "2.  General Counsel deploys whistleblower mechanism within 60 days.",
        "3.  CEO commits to SBTi steel pathway and issues public statement.",
        "4.  CFO/CSO appoints CSRD programme manager and completes DMA by Q4 2024.",
        "5.  ESGIntel provides quarterly monitoring and investor-grade progress reporting.",
    ]
    sy = 3.42
    for step in steps:
        _add_text(s, step, 0.7, sy, 11.5, 0.42,
                  size=10, color=WHITE, font=D.BRAND["font_body"])
        sy += 0.50

    _add_text(s, D.BRAND["platform"] + "  |  " + D.BRAND["tagline"],
              0.7, 6.55, 10.0, 0.45,
              size=9, color=MUTED, font=D.BRAND["font_body"])
    _add_text(s, D.BRAND["confidential"],
              0.7, 7.05, 11.5, 0.35,
              size=7.5, color=MUTED, font=D.BRAND["font_body"])

    prs.save(out_path)


# ---------------------------------------------------------------------------
# ═══════════════════════════════════════════════════════════════
# INVESTOR DECK  (~8 slides)
# ═══════════════════════════════════════════════════════════════
# ---------------------------------------------------------------------------

def _build_investor_deck(tmpdir: str, out_path: str):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]

    today = datetime.date.today().strftime("%B %Y")

    # ── Slide 1: Cover ────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    _add_rect(s, 0, 0, 13.333, 7.5, D.BRAND["ink"])
    _add_rect(s, 0, 0, 0.42, 7.5, D.BRAND["accent"])
    _add_rect(s, 0.42, 4.2, 12.91, 0.04, D.BRAND["accent"])

    _add_text(s, "ESGIntel", 0.7, 0.55, 11.0, 0.6,
              size=11, color=ACCENT, font=D.BRAND["font_body"])
    _add_text(s, D.COMPANY["name"], 0.7, 1.15, 11.0, 1.1,
              size=40, bold=True, color=WHITE, font=D.BRAND["font_head"])
    _add_text(s, "Investor Presentation — ESG Due Diligence", 0.7, 2.35, 11.0, 0.7,
              size=20, bold=False, color=ACCENT, font=D.BRAND["font_head"], italic=True)
    _add_text(s, D.BRAND["confidential"], 0.7, 3.10, 11.0, 0.4,
              size=9, color=MUTED, font=D.BRAND["font_body"])
    _add_text(s, today, 0.7, 3.6, 3.5, 0.35,
              size=9, color=MUTED, font=D.BRAND["font_body"])

    meta_y = 4.50
    for label, val in [
        ("Sector",    D.COMPANY["sector"]),
        ("Revenue",   f"€{D.COMPANY['revenue_eur_bn']}bn  |  {D.COMPANY['employees']:,} employees"),
        ("Domicile",  D.COMPANY["domicile"]),
        ("ESG Score", "38/100 — High Risk — P19 vs NACE C24"),
    ]:
        _add_text(s, label.upper(), 0.7, meta_y, 2.0, 0.25,
                  size=7, bold=True, color=ACCENT, font=D.BRAND["font_body"])
        _add_text(s, val, 2.6, meta_y, 9.5, 0.28,
                  size=8.5, color=WHITE, font=D.BRAND["font_body"])
        meta_y += 0.36
    _add_text(s, "CONFIDENTIAL", 9.5, 7.1, 3.5, 0.35,
              size=7.5, bold=True, color=MUTED, align=PP_ALIGN.RIGHT,
              font=D.BRAND["font_body"])

    # ── Slide 2: ESG Score + Peer Position ────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    _add_chrome(s, "ESG Score & Peer Positioning", 2)
    _action_headline(s, "38/100 High Risk (P19) — bottom quartile across all three pillars vs 47 NACE C24 peers")

    # Score tile
    _big_number_tile(s, "38 / 100", "Overall ESG Risk Score",
                     MARGIN_L, TITLE_BAR_H + 0.72, 2.9, 1.55,
                     bg_hex=D.BRAND["red"], val_color="FFFFFF")
    _add_text(s, "P19 vs peers — bottom quartile\nHigh Risk label\n72% data confidence",
              MARGIN_L + 3.05, TITLE_BAR_H + 0.90, 4.5, 1.10,
              size=11, bold=False, color=SLATE, font=D.BRAND["font_body"])

    # Pillar summary tiles
    tile_x = MARGIN_L
    for pillar, data in D.SCORES["pillars"].items():
        _big_number_tile(s, str(data["score"]) + "/100", f"{pillar}\n{data['label']} · P{data['percentile']}",
                         tile_x, TITLE_BAR_H + 2.55, 3.95, 1.25,
                         bg_hex=rag_hex(data["label"]), val_color="FFFFFF")
        tile_x += 4.1

    # Pillar chart
    chart_path = _chart_pillar_scores(tmpdir)
    s.shapes.add_picture(chart_path, Inches(MARGIN_L + 7.5), Inches(TITLE_BAR_H + 0.72),
                         Inches(5.0), Inches(2.5))

    # Investor framing
    frame_y = TITLE_BAR_H + 4.00
    _add_rect(s, MARGIN_L, frame_y, CONTENT_W, 1.45, D.BRAND["band_bg"])
    _add_text(s, "Investor Framing:",
              MARGIN_L + 0.15, frame_y + 0.10, CONTENT_W - 0.3, 0.28,
              size=9, bold=True, color=INK, font=D.BRAND["font_head"])
    _add_text(s,
              f"A 38/100 High-Risk score signals that {D.COMPANY['short']}'s ESG profile currently limits its access to green finance, "
              "increases its cost of carbon compliance, and exposes it to regulatory penalties. The P19 peer ranking means 81% "
              "of comparable steel companies present lower ESG risk. Remediation of the engagement asks would likely improve the "
              "score by 10–15 points, unlocking the green-bond market and reducing ETS exposure through a credible decarbonisation plan.",
              MARGIN_L + 0.15, frame_y + 0.42, CONTENT_W - 0.3, 0.95,
              size=8.5, color=SLATE, font=D.BRAND["font_body"])

    # ── Slide 3: Top Risks ────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    _add_chrome(s, "Top ESG Risks", 3)
    _action_headline(s, "Three critical-scored risks account for the majority of quantifiable financial exposure")

    top_risks = sorted(D.RISKS, key=lambda x: x["score"], reverse=True)[:6]
    for i, r in enumerate(top_risks):
        rx = MARGIN_L + (i % 2) * 6.3
        ry = TITLE_BAR_H + 0.72 + (i // 2) * 1.7
        sev = D.severity(r["score"])
        _add_rect(s, rx, ry, 6.05, 1.55, D.BRAND["band_bg"])
        _add_rect(s, rx, ry, 6.05, 0.05, rag_hex(sev))
        _rag_chip(s, f"{sev}  {r['score']}", sev, rx + 0.12, ry + 0.10, w=1.55, h=0.28)
        _add_text(s, r["id"], rx + 1.80, ry + 0.10, 0.72, 0.28,
                  size=8, bold=True, color=rgb(rag_hex(sev)), font=D.BRAND["font_body"])
        _add_text(s, r["title"], rx + 2.55, ry + 0.08, 3.35, 0.32,
                  size=8.5, bold=True, color=INK, font=D.BRAND["font_head"])
        _add_text(s, r["financial"], rx + 0.12, ry + 0.50, 2.8, 0.30,
                  size=8.5, bold=True, color=rgb(rag_hex(sev)), font=D.BRAND["font_body"])
        _add_text(s, r["summary"][:120] + "…", rx + 0.12, ry + 0.85, 5.8, 0.60,
                  size=7.5, color=SLATE, font=D.BRAND["font_body"])

    # ── Slide 4: Climate Exposure ─────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    _add_chrome(s, "Climate Financial Exposure", 4)
    _action_headline(s, "ETS (€68M) + CBAM (€24M) = €92M recurring annual exposure by 2027-30, with no hedging in place")

    # KPI tiles for ETS and CBAM
    _big_number_tile(s, "€68M/yr", "EU ETS Net Cost by 2030",
                     MARGIN_L, TITLE_BAR_H + 0.72, 3.8, 1.5,
                     bg_hex=D.BRAND["red"], val_color="FFFFFF")
    _big_number_tile(s, "€24M/yr", "CBAM Certificates by 2027",
                     MARGIN_L + 3.95, TITLE_BAR_H + 0.72, 3.8, 1.5,
                     bg_hex=D.BRAND["amber"], val_color="FFFFFF")
    _big_number_tile(s, "€92M", "Combined Annual Run-Rate",
                     MARGIN_L + 7.90, TITLE_BAR_H + 0.72, 4.35, 1.5,
                     bg_hex=D.BRAND["ink"], val_color="FFFFFF")

    # Climate scenarios chart
    chart_path = _chart_climate_scenarios(tmpdir)
    s.shapes.add_picture(chart_path, Inches(MARGIN_L), Inches(TITLE_BAR_H + 2.45),
                         Inches(6.8), Inches(3.0))

    # Commentary
    comm_x = MARGIN_L + 7.0
    _add_rect(s, comm_x, TITLE_BAR_H + 2.45, CONTENT_W - 7.0, 3.0, D.BRAND["band_bg"])
    _add_text(s, "Key Drivers", comm_x + 0.15, TITLE_BAR_H + 2.55, CONTENT_W - 7.2, 0.30,
              size=9, bold=True, color=INK, font=D.BRAND["font_head"])
    bullets = [
        "Scope 1: 842 ktCO2e — BF-BOF locked",
        "Free allocation: ~40% → ~10% by 2030",
        "No ETS hedging or internal carbon price",
        "CBAM: ~€280M non-EU revenue exposed",
        "No SBTi commitment — no carbon roadmap",
        "DRI/EAF transition capex: €400–600M",
    ]
    for bi, b in enumerate(bullets):
        _add_text(s, "•  " + b, comm_x + 0.15, TITLE_BAR_H + 2.92 + bi * 0.38,
                  CONTENT_W - 7.25, 0.32,
                  size=8, color=SLATE, font=D.BRAND["font_body"])

    # ── Slide 5: Financial Exposure Chart ─────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    _add_chrome(s, "Quantified Financial Exposure", 5)
    _action_headline(s, "Total quantified ESG exposure exceeds €950M — ETS, stranding and CSDDD are the three largest items")

    chart_path = _chart_financial_exposure(tmpdir)
    s.shapes.add_picture(chart_path, Inches(MARGIN_L), Inches(TITLE_BAR_H + 0.68),
                         Inches(7.0), Inches(3.7))

    _add_table(s, ["Driver", "€M", "Type", "Rating"],
               [[f["driver"], f"€{f['amount_eur_m']}M", f["type"], f["rating"]]
                for f in D.FINANCIAL_EXPOSURE],
               MARGIN_L + 7.2, TITLE_BAR_H + 0.68, CONTENT_W - 7.2, 3.7,
               col_widths=[2.5, 0.72, 1.3, 0.8], font_size=8)

    total = sum(f["amount_eur_m"] for f in D.FINANCIAL_EXPOSURE)
    _big_number_tile(s, f"€{total}M", "Total Quantified ESG Exposure",
                     MARGIN_L, TITLE_BAR_H + 4.55, 7.0, 0.80,
                     bg_hex=D.BRAND["red"], val_color="FFFFFF")

    # ── Slide 6: Strategic Priorities / Transition ─────────────────────────
    s = prs.slides.add_slide(blank_layout)
    _add_chrome(s, "Strategic Priorities — ESG Transition Roadmap", 6)
    _action_headline(s, "Three horizons define the remediation path — near-term wins unlock medium-term financing and long-term value")

    horizons = [
        ("Immediate (0–12 months)", D.BRAND["red"],
         ["Deploy whistleblower mechanism (60 days — non-negotiable)",
          "Appoint CSRD programme manager; launch DMA",
          "Commit to SBTi and publish public statement",
          "Register in EU CBAM portal; appoint compliance officer",
          "Add ESG KPIs (LTIFR, GHG) to executive LTIP"]),
        ("Medium-Term (12–36 months)", D.BRAND["amber"],
         ["Publish ESRS E1 transition plan (DRI/EAF roadmap by 2025)",
          "Complete FY2025 CSRD report (first limited assurance)",
          "Achieve LTIFR ≤1.5 and ISO 45001 at all sites",
          "Invest €12–18M in closed-loop water recycling",
          "Launch EU Green Bond or SLL (contingent on SBTi)"]),
        ("Long-Term (36–120 months)", D.BRAND["accent"],
         ["Commission DRI/EAF at Taranto; phased BF-BOF exit 2032+",
          "Target 50% GHG intensity reduction vs 2023 baseline",
          "Secure EU Innovation Fund grant (€50–120M)",
          "Achieve investment-grade ESG rating (>55/100, P40+)",
          "Automotive Scope 3 supplier approval (€456M revenue)"]),
    ]
    for i, (title, color_h, bullets) in enumerate(horizons):
        hx = MARGIN_L + i * 4.18
        hy = TITLE_BAR_H + 0.72
        _add_rect(s, hx, hy, 4.0, 5.5, D.BRAND["band_bg"])
        _add_rect(s, hx, hy, 4.0, 0.07, color_h)
        _add_text(s, title, hx + 0.12, hy + 0.12, 3.76, 0.42,
                  size=10, bold=True, color=INK, font=D.BRAND["font_head"])
        for bi, b in enumerate(bullets):
            _add_text(s, "•  " + b, hx + 0.12, hy + 0.65 + bi * 0.90,
                      3.76, 0.80,
                      size=8.5, color=SLATE, font=D.BRAND["font_body"])

    # ── Slide 7: Green Finance Upside ─────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    _add_chrome(s, "Green Finance Upside", 7)
    _action_headline(s, "€450–600M of lower-cost ESG financing is accessible — contingent on SBTi target and CSRD readiness")

    for i, opp in enumerate(D.OPPORTUNITIES):
        cx = MARGIN_L + i * 4.18
        cy = TITLE_BAR_H + 0.72
        _add_rect(s, cx, cy, 4.0, 4.85, D.BRAND["band_bg"])
        _add_rect(s, cx, cy, 4.0, 0.08, D.BRAND["accent"])
        _add_text(s, opp["lever"], cx + 0.15, cy + 0.14, 3.7, 0.52,
                  size=11, bold=True, color=INK, font=D.BRAND["font_head"])
        _big_number_tile(s, opp["size"], "Financing Size",
                         cx + 0.15, cy + 0.78, 3.7, 1.0,
                         bg_hex=D.BRAND["accent_lt"], val_color=D.BRAND["ink"])
        _add_text(s, "Benefit:", cx + 0.15, cy + 1.88, 3.7, 0.25,
                  size=8, bold=True, color=INK, font=D.BRAND["font_body"])
        _add_text(s, opp["benefit"], cx + 0.15, cy + 2.15, 3.7, 0.72,
                  size=8.5, color=SLATE, font=D.BRAND["font_body"])
        _add_rect(s, cx + 0.15, cy + 2.95, 3.7, 0.03, D.BRAND["rule"])
        _add_text(s, "Pre-condition:", cx + 0.15, cy + 3.05, 3.7, 0.25,
                  size=8, bold=True, color=AMBER, font=D.BRAND["font_body"])
        _add_text(s, opp["precondition"], cx + 0.15, cy + 3.32, 3.7, 1.0,
                  size=8, color=SLATE, font=D.BRAND["font_body"])

    # Total upside callout
    _big_number_tile(s, "€450–600M", "Total Green Finance Upside",
                     MARGIN_L + 12.55 - 3.0, TITLE_BAR_H + 0.72, 2.8, 1.65,
                     bg_hex=D.BRAND["accent"], val_color="FFFFFF")

    # ── Slide 8: Engagement Asks / Close ──────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    _add_chrome(s, "Engagement Asks & Next Steps", 8)
    _action_headline(s, "Six engagement asks — investor is seeking formal Board commitment to these actions within 30 days")

    _add_table(s, ["Priority", "Engagement Ask", "Owner", "Timeline"],
               [[ea["priority"], ea["ask"], ea["owner"], ea["by"]]
                for ea in D.ENGAGEMENT_ASKS],
               MARGIN_L, TITLE_BAR_H + 0.72, CONTENT_W, 4.25,
               col_widths=[0.88, 6.2, 2.2, 1.4], font_size=9)

    # Closing statement
    close_y = TITLE_BAR_H + 5.15
    _add_rect(s, MARGIN_L, close_y, CONTENT_W, 1.28, D.BRAND["ink"])
    _add_text(s, "ESGIntel is prepared to provide quarterly monitoring updates, investor-grade progress reporting, "
                 f"and co-ordination with {D.COMPANY['short']} management on all engagement asks. "
                 "Contact: " + D.BRAND["platform"] + "  |  " + D.BRAND["tagline"],
              MARGIN_L + 0.25, close_y + 0.18, CONTENT_W - 0.5, 0.95,
              size=9.5, color=WHITE, font=D.BRAND["font_body"])

    prs.save(out_path)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def build_deck(deck_key: str, out_dir: str) -> str:
    """
    Generate a McKinsey-style ESG engagement deck.

    Parameters
    ----------
    deck_key : str
        'board'    → ESG Engagement Deck (Board Briefing), ~15 slides
        'investor' → Investor Presentation (ESG), ~8 slides
    out_dir : str
        Directory where the .pptx file will be saved.

    Returns
    -------
    str
        Absolute path to the generated .pptx file.
    """
    if deck_key not in ("board", "investor"):
        raise ValueError(f"deck_key must be 'board' or 'investor', got {deck_key!r}")

    os.makedirs(out_dir, exist_ok=True)

    filenames = {
        "board":    "ESG_Engagement_Deck_Board.pptx",
        "investor": "Investor_Presentation_ESG.pptx",
    }
    out_path = os.path.abspath(os.path.join(out_dir, filenames[deck_key]))

    with tempfile.TemporaryDirectory() as tmpdir:
        if deck_key == "board":
            _build_board_deck(tmpdir, out_path)
        else:
            _build_investor_deck(tmpdir, out_path)

    return out_path


# ---------------------------------------------------------------------------
# CLI entry-point for quick testing
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    import sys
    out_dir = sys.argv[1] if len(sys.argv) > 1 else "/tmp"
    for key in ("board", "investor"):
        path = build_deck(key, out_dir)
        size_kb = os.path.getsize(path) // 1024
        print(f"[{key}]  {path}  ({size_kb} KB)")
